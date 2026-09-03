# -*- coding: utf-8 -*-
"""Fill the SIH 2026 idea-submission template for HeatLens.

Every figure quoted here comes from the repository's own outputs
(web/data/*.json) or from scripts/05_kill_gate.py.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from deckkit import *

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
SRC = os.path.join(HERE, "template.pptx")
OUT = os.path.join(HERE, "HeatLens-SIH2026-Idea.pptx")

TEAM = "HeatLens"


ICO = os.path.join(HERE, "icons")


def img(name):
    return os.path.join(IMG, name)


def picto(s, name, x, y, size):
    """Drop one of the rendered react-icons onto a slide."""
    return s.shapes.add_picture(os.path.join(ICO, name + ".png"),
                                E(x), E(y), height=E(size))


# =========================================================== slide 1: title
def slide1(s):
    sub_ph = by_name(s, "Subtitle 3")
    sub_ph.left, sub_ph.top = E(0.36), E(1.28)
    sub_ph.width, sub_ph.height = E(6.95), E(0.80)
    stf = sub_ph.text_frame
    stf.word_wrap = True
    stf.vertical_anchor = MSO_ANCHOR.BOTTOM
    stf.margin_left = stf.margin_right = 0
    r = retext(sub_ph, "Heat warnings measured per neighbourhood, "
                       "per body, per hour")
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.name = HEAD
    r.font.color.rgb = INK
    stf.paragraphs[0].alignment = PP_ALIGN.LEFT

    tb = by_name(s, "TextBox 9")
    tb.left, tb.top, tb.width, tb.height = E(0.36), E(2.34), E(6.95), E(4.45)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for r in list(tf.paragraphs[0].runs):
        r._r.getparent().remove(r._r)

    fields = [
        ("Problem Statement ID –", "26083", 20),
        ("Problem Statement Title-",
         "Extreme Heatwave Early Warning and\nHuman Thermal Stress Index", 17),
        ("Theme-", "Disaster Management", 18),
        ("PS Category- Software", None, 18),
        ("Team Name (Registered on portal)", TEAM, 20),
    ]
    first = True
    for label, value, vsize in fields:
        p = para(tf, first=first, space_before=0 if first else 11, line=0.98)
        run(p, label, 11.5, bold=True, color=MUTED, font="Arial")
        first = False
        if value is None:
            continue
        for i, line in enumerate(value.split("\n")):
            q = para(tf, space_before=2 if i == 0 else 0, line=0.94)
            run(q, line, vsize, bold=True, color=INK, font="Arial")

    p = para(tf, space_before=20, line=1.12)
    run(p, "A heatwave warning tells you the air temperature. It does not "
           "tell you what that heat does to a person.", 12, italic=True,
        color=SLATE)
    p = para(tf, space_before=5, line=1.12)
    run(p, "Prototype: 392 neighbourhoods · 135 passing physics tests · "
           "runs offline", 11, bold=True, color=CRIMSON)


# ==================================================== slide 2: idea / solution
def slide2(s):
    """Plain language only. The picture carries the explanation."""
    t = by_name(s, "Title 1")
    t.left, t.top, t.width, t.height = E(1.86), E(0.20), E(2.40), E(0.80)
    r = retext(t, "HeatLens")
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.name = HEAD
    r.font.color.rgb = INK
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    tag = txbox(s, 4.00, 0.34, 6.50, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tag.text_frame, first=True, line=1.05)
    run(p, "A heatwave warning measured on the body, not the thermometer",
        14.5, color=CRIMSON)

    drop(by_name(s, "TextBox 8"))

    pointer(s, L, 1.04, 12.49,
            "Proposed Solution (Describe your Idea/Solution/Prototype)", 12.5)

    picture(s, img("d2_story.png"), L, 1.36, w=12.49)

    cols = [
        ("Detailed explanation of the proposed solution", [
            "We cut the city into 392 small zones, each about the size of a "
            "neighbourhood.",
            "For every zone and every hour we work out what the heat does to a "
            "body — using the limits doctors and safety inspectors already use.",
            "That becomes one plain number: how many minutes it is safe to be "
            "outside.",
            "The whole thing is a single file that opens with no internet.",
        ]),
        ("How it addresses the problem", [
            "Today a whole district gets one temperature, one or two days "
            "ahead.",
            "That number cannot tell a dry 40 °C from a sticky 40 °C — and only "
            "one of those is workable outdoors.",
            "We answer street by street, five days ahead, in the language "
            "people actually speak.",
        ]),
        ("Innovation and uniqueness of the solution", [
            "We pick the right measure for the weather that day. In May 2010 "
            "the usual worker index hid the danger; the “feels like” one "
            "showed it.",
            "We count the nights. Seven in a row never dropped below 27 °C — no "
            "daytime warning can see that.",
            "We print what we do not know on the face of the tool.",
        ]),
    ]
    cw, gap = 4.03, 0.20
    for i, (head, items) in enumerate(cols):
        x = L + i * (cw + gap)
        sub(s, x, 4.80, cw, head)
        body(s, x, 5.10, cw, 1.70, items, 9.2)

# ================================================= slide 3: technical approach
def slide3(s):
    """The whole project in one picture: what goes in, what happens, what
    comes out. Technical names ride along as small secondary text."""
    section_title(s, "TECHNICAL APPROACH")
    drop(by_name(s, "TextBox 8"))

    pointer(s, L, 1.00, 8.0,
            "Technologies to be used (e.g. programming languages, frameworks, "
            "hardware)", 11.5)

    # Labels stay short so every chip holds its value on one line at 3.02" wide
    stack = [
        ("CORE", "Python 3.12 · NumPy · PyYAML"),
        ("PHYSICS", "thermofeel (ECMWF) · Liljegren · UTCI"),
        ("PHYSIOLOGY", "ISO 7243 · ACGIH TLV · ISO 7726"),
        ("GEOSPATIAL", "Uber H3 v4 · OpenStreetMap / Overpass"),
        ("WEATHER", "Open-Meteo ERA5 archive + forecast"),
        ("FRONTEND", "React 18 · TypeScript · Vite · SVG map"),
        ("DELIVERY", "one inlined index.html · CAP 1.2 XML"),
        ("QUALITY", "pytest ×135 · GitHub Actions, 6-hourly"),
    ]
    cw, gp, ch = 3.02, 0.135, 0.31
    for i, (lab, value) in enumerate(stack):
        x = L + (i % 4) * (cw + gp)
        y = 1.30 + (i // 4) * (ch + 0.035)
        card(s, x, y, cw, ch, fill=PANEL, line=BORDER, radius=0.16)
        tb = txbox(s, x + 0.11, y + 0.02, cw - 0.22, ch - 0.04,
                   anchor=MSO_ANCHOR.MIDDLE)
        p = para(tb.text_frame, first=True, line=1.0)
        run(p, lab + "  ", 7.6, bold=True, color=MAGENTA)
        run(p, value, 8.4, color=SLATE)

    pointer(s, L, 2.05, 12.49,
            "Methodology and process for implementation (Flow Charts/Images/ "
            "working prototype)", 11.5)

    picture(s, img("d3_system.png"), L, 2.34, w=12.49)

# ============================================== slide 4: feasibility/viability
def slide4(s):
    """A decision flowchart on the left, gap-to-fix arrows on the right."""
    section_title(s, "FEASIBILITY AND VIABILITY")
    drop(by_name(s, "TextBox 8"))

    lw, rx, rw = 6.12, 6.88, 6.03

    pointer(s, L, 1.02, lw, "Analysis of the feasibility of the idea", 11.5)
    cap = txbox(s, L + 0.20, 1.28, lw - 0.20, 0.22)
    p = para(cap.text_frame, first=True, line=1.0)
    run(p, "Before building the product we tested whether the idea was true.",
        8.6, color=MUTED)

    pointer(s, rx, 1.02, rw, "Potential challenges and risks", 11.5,
            color=CRIMSON)
    pointer(s, rx, 1.26, rw, "Strategies for overcoming these challenges",
            11.5, color=GREEN)

    picture(s, img("d4_gate.png"), L, 1.50, w=lw)
    picture(s, img("d4_risks.png"), rx, 1.50, w=rw)


# ================================================ slide 5: impact and benefits
def slide5(s):
    """One computed hour, read five ways — then the two charts that prove it."""
    section_title(s, "IMPACT AND BENEFITS")
    drop(by_name(s, "TextBox 8"))

    lw, rx, rw = 7.70, 8.42, 4.49

    pointer(s, L, 1.02, lw, "Potential impact on the target audience", 11.5)
    picture(s, img("d5_people.png"), L, 1.32, w=lw)
    picture(s, img("v_worksafety.png"), L, 4.16, w=lw)
    cap = txbox(s, L, 6.42, lw, 0.34)
    p = para(cap.text_frame, first=True, line=1.05)
    run(p, "Eight hours straight, on 21 May 2010, when no one could safely "
           "work outdoors at all.", 8.8, color=MUTED)

    pointer(s, rx, 1.02, rw,
            "Benefits of the solution (social, economic, environmental, etc.)",
            11.5, h=0.50)

    bens = [
        ("users", "Social", "Protects people who cannot simply stay indoors.",
         CRIMSON),
        ("rupee", "Economic",
         "Moving hours saves work a blanket stop-work order loses.", EMBER),
        ("leaf", "Environmental",
         "Shows where trees and cool roofs buy the most.", GREEN),
        ("file", "Institutional",
         "Speaks the national alert format already in use.", STEEL),
    ]
    y, bh, bg = 1.60, 0.56, 0.07
    for i, (ic, name, what, col) in enumerate(bens):
        yy = y + i * (bh + bg)
        card(s, rx, yy, rw, bh, fill=WHITE, line=BORDER)
        picto(s, ic, rx + 0.14, yy + 0.15, 0.26)
        tb = txbox(s, rx + 0.50, yy + 0.08, rw - 0.64, bh - 0.16)
        p = para(tb.text_frame, first=True, line=1.0)
        run(p, name, 9.4, bold=True, color=col)
        p = para(tb.text_frame, line=1.03)
        run(p, what, 8.0, color=SLATE)

    picture(s, img("v_nights.png"), rx, 4.16, w=rw)
    cap = txbox(s, rx, 6.53, rw, 0.30)
    p = para(cap.text_frame, first=True, line=1.05)
    run(p, "The part of a heatwave a daytime warning never sees.", 8.8,
        color=MUTED)

# ============================================ slide 6: research and references
def slide6(s):
    section_title(s, "RESEARCH AND REFERENCES")
    drop(by_name(s, "TextBox 8"))

    pointer(s, L, 1.10, 12.49,
            "Details / Links of the reference and research work", 12.5)

    cols = [
        ("Physics & standards implemented", CRIMSON, [
            ("Liljegren et al. (2008), J. Occup. Environ. Hyg.",
             "Solar-corrected WBGT; used via thermofeel (ECMWF)."),
            ("Bröde et al. (2012), Int. J. Biometeorol.",
             "The operational UTCI polynomial."),
            ("Stull (2011), J. Appl. Meteorol. Climatol.",
             "doi:10.1175/JAMC-D-11-0143.1 — wet-bulb from T and RH."),
            ("Rothfusz (1990) / NOAA NWS", "Heat Index regression and "
                                           "adjustments."),
            ("ISO 7243:2017 · ISO 7726 · ISO 8996",
             "WBGT limits, mean radiant temperature, metabolic rate."),
            ("ACGIH TLV®", "Work/rest allocation tables per metabolic class."),
        ]),
        ("Data sources — free, global, no institutional access", STEEL, [
            ("Open-Meteo", "ERA5 hourly archive and forecast API: temperature, "
                           "humidity, wind, GHI/DNI/DHI."),
            ("OpenStreetMap via Overpass API",
             "Road, green and water cover per H3 cell."),
            ("Uber H3 v4", "Resolution-8 discrete global grid, ~0.74 km² cells."),
            ("thermofeel (ECMWF)",
             "Operational reference implementation we validate against."),
            ("OASIS CAP 1.2",
             "Alert interchange format; NDMA SACHET compatible."),
        ]),
        ("Evidence base for the problem", EMBER, [
            ("Azhar et al. (2014), PLoS ONE",
             "Excess mortality in the May 2010 Ahmedabad heatwave."),
            ("Knowlton et al. (2014), IJERPH",
             "Design of South Asia's first Heat Action Plan."),
            ("Hess et al. (2018)",
             "Pilot evaluation of the Ahmedabad Heat Action Plan."),
            ("IMD heatwave criteria",
             "The district-scale, dry-bulb baseline this improves on."),
            ("NDMA, National Guidelines: Heat Wave (2019)",
             "The action framework our outputs are shaped to feed."),
            ("Primary-source verification", "Casualty figures are checked "
                                            "against the source papers before "
                                            "any are quoted on stage."),
        ]),
    ]
    cw, gap = 4.03, 0.20
    for i, (head, col, items) in enumerate(cols):
        x = L + i * (cw + gap)
        tb = txbox(s, x, 1.48, cw, 0.30)
        p = para(tb.text_frame, first=True, line=1.02)
        run(p, head, 10.2, bold=True, color=col)
        y = 1.84
        for title, note in items:
            hexmark(s, x, y + 0.04, 0.095, col)
            t2 = txbox(s, x + 0.17, y - 0.025, cw - 0.17, 0.52)
            p = para(t2.text_frame, first=True, line=1.02)
            run(p, title, 8.6, bold=True, color=INK)
            p = para(t2.text_frame, line=1.04)
            run(p, note, 8.2, color=MUTED)
            y += 0.53

    # ---------------- provenance band
    tb = txbox(s, L, 5.10, 12.49, 0.30)
    p = para(tb.text_frame, first=True, line=1.02)
    run(p, "Provenance of every layer we ship — ", 10.4, bold=True, color=INK)
    run(p, "stated in the UI and in web/data/meta.json, because a number "
           "without its status is not evidence.", 9.4, color=MUTED)

    layers = [
        ("Past weather", "Open-Meteo ERA5", "MEASURED", GREEN),
        ("Urban form", "OpenStreetMap", "MEASURED", GREEN),
        ("Heat-stress indices", "thermofeel / ECMWF", "MEASURED", GREEN),
        ("Safe limits for a body", "ISO 7243 + ACGIH", "PUBLISHED STANDARD",
         STEEL),
        ("Urban-heat amplitude", "literature, 3.0 °C", "ASSUMED", EMBER),
        ("Vulnerability", "city-wide constant", "NOT FITTED", CRIMSON),
        ("Health risk", "literature-shaped", "NOT CALIBRATED", CRIMSON),
    ]
    cw, gap = 1.66, 0.135
    for i, (name, src, status, col) in enumerate(layers):
        x = L + i * (cw + gap)
        card(s, x, 5.48, cw, 1.16, fill=WHITE, line=BORDER)
        tb = txbox(s, x + 0.11, 5.58, cw - 0.22, 0.62)
        p = para(tb.text_frame, first=True, line=1.02)
        run(p, name, 9.0, bold=True, color=INK)
        p = para(tb.text_frame, space_before=2, line=1.02)
        run(p, src, 7.8, color=MUTED)
        pill = card(s, x + 0.11, 6.24, cw - 0.22, 0.28, fill=col, line=None,
                    radius=0.42)
        tfp = pill.text_frame
        tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = para(tfp, first=True, align=PP_ALIGN.CENTER, line=0.95)
        run(p, status, 6.8, bold=True, color=WHITE)


# ---------------------------------------------------------------- shared bits
def section_title(s, text):
    t = by_name(s, "Title 1")
    t.left, t.top, t.width, t.height = E(1.86), E(0.20), E(8.70), E(0.80)
    r = retext(t, text)
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.name = HEAD
    r.font.color.rgb = INK


def sub(s, x, y, w, text):
    """A required sub-pointer, rendered verbatim as a sub-heading."""
    tb = txbox(s, x, y, w, 0.26)
    p = para(tb.text_frame, first=True, line=1.0)
    run(p, text, 10.8, bold=True, color=INK)
    return tb


def body(s, x, y, w, h, items, size):
    tb = txbox(s, x, y, w, h)
    tf = tb.text_frame
    for i, item in enumerate(items):
        p = para(tf, first=(i == 0), space_before=0 if i == 0 else 5, line=1.03)
        bullet(p)
        if isinstance(item, tuple):
            run(p, item[0], size, bold=True, color=INK)
            run(p, item[1], size, color=SLATE)
        else:
            run(p, item, size, color=SLATE)
    return tb


def main():
    prs = Presentation(SRC)
    delete_slide(prs, 6)                     # the template's instructions slide
    sl = prs.slides

    for s in list(sl)[1:]:                   # team badge on every content slide
        for nm in ("Oval 9", "Oval 10", "Oval 11", "Oval 8"):
            o = by_name(s, nm)
            if o is not None:
                # Shrink the logo placeholder so it stops crowding the first
                # pointer heading on every content slide.
                o.left, o.top = E(0.36), E(0.20)
                o.width, o.height = E(1.24), E(0.64)
                r = retext(o, TEAM)
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.color.rgb = INK
                break

    slide1(sl[0])
    slide2(sl[1])
    slide3(sl[2])
    slide4(sl[3])
    slide5(sl[4])
    slide6(sl[5])
    prs.save(OUT)
    print("wrote", OUT, os.path.getsize(OUT) // 1024, "KB")


if __name__ == "__main__":
    main()
