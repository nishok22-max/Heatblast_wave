# -*- coding: utf-8 -*-
"""Layout helpers for filling the SIH idea-submission template."""
import copy
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---- palette -------------------------------------------------------------
INK      = RGBColor(0x24, 0x10, 0x43)
SLATE    = RGBColor(0x3A, 0x3F, 0x52)
MUTED    = RGBColor(0x7A, 0x7F, 0x91)
CRIMSON  = RGBColor(0x8E, 0x10, 0x38)
MAGENTA  = RGBColor(0xC2, 0x18, 0x5B)
EMBER    = RGBColor(0xD2, 0x54, 0x1F)
AMBER    = RGBColor(0xF5, 0xA0, 0x3C)
GREEN    = RGBColor(0x2E, 0x7D, 0x5B)
STEEL    = RGBColor(0x1B, 0x6C, 0xA8)
PANEL    = RGBColor(0xF5, 0xF2, 0xF8)
SAND     = RGBColor(0xFD, 0xF4, 0xE8)
MIST     = RGBColor(0xEE, 0xF4, 0xF9)
BORDER   = RGBColor(0xE1, 0xDC, 0xEA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

BODY = "Calibri"
HEAD = "Times New Roman"

# ---- geometry ------------------------------------------------------------
L, R = 0.42, 12.91
TOP, BOT = 1.14, 6.82


def E(v):
    return Inches(v)


# ---- primitives ----------------------------------------------------------
def drop(shape):
    shape._element.getparent().remove(shape._element)


def by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def txbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].text = ""
    return tb


def para(tf, first=False, space_before=0, space_after=0, line=1.0,
         align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p


_RPR_ORDER = ["a:ln", "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
              "a:pattFill", "a:grpFill", "a:effectLst", "a:effectDag",
              "a:highlight", "a:uLnTx", "a:uLn", "a:uFillTx", "a:uFill",
              "a:latin", "a:ea", "a:cs", "a:sym", "a:hlinkClick",
              "a:hlinkMouseOver", "a:rtl", "a:extLst"]


def _reorder_rpr(rPr):
    idx = {qn(t): i for i, t in enumerate(_RPR_ORDER)}
    for k in sorted(list(rPr), key=lambda e: idx.get(e.tag, 99)):
        rPr.append(k)


def run(p, text, size, bold=False, color=SLATE, font=BODY, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    # Without an explicit ea/cs typeface PowerPoint routes ambiguous-width
    # characters (° × –) to the theme's CJK font and renders them full-width.
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font)
    _reorder_rpr(rPr)
    return r


# a:pPr children must appear in schema order or PowerPoint rejects the file
_PPR_ORDER = ["a:lnSpc", "a:spcBef", "a:spcAft", "a:buClrTx", "a:buClr",
              "a:buSzTx", "a:buSzPct", "a:buSzPts", "a:buFontTx", "a:buFont",
              "a:buNone", "a:buAutoNum", "a:buChar", "a:tabLst", "a:defRPr",
              "a:extLst"]


def _reorder_ppr(pPr):
    idx = {qn(t): i for i, t in enumerate(_PPR_ORDER)}
    kids = list(pPr)
    kids.sort(key=lambda e: idx.get(e.tag, 99))
    for k in kids:
        pPr.append(k)


def bullet(p, char="•", color=CRIMSON, mar=0.115, size_pct=95):
    """Give one paragraph a real hanging bullet."""
    pPr = p._p.get_or_add_pPr()
    marL = int(Inches(mar))
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-marL))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr",
                "a:buFont", "a:buSzPct"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    clr = etree.SubElement(buClr, qn("a:srgbClr"))
    clr.set("val", str(color))
    sz = etree.SubElement(pPr, qn("a:buSzPct"))
    sz.set("val", str(size_pct * 1000))
    bf = etree.SubElement(pPr, qn("a:buFont"))
    bf.set("typeface", "Arial")
    bc = etree.SubElement(pPr, qn("a:buChar"))
    bc.set("char", char)
    _reorder_ppr(pPr)
    return p


def shape(slide, kind, x, y, w, h, fill=None, line=None, lw=0.75, radius=None):
    sp = slide.shapes.add_shape(kind, E(x), E(y), E(w), E(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    sp.text_frame.word_wrap = True
    sp.text_frame.margin_left = sp.text_frame.margin_right = 0
    sp.text_frame.margin_top = sp.text_frame.margin_bottom = 0
    return sp


def card(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=0.055):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                 fill=fill, line=line, radius=radius)


def hexmark(slide, x, y, size=0.125, color=CRIMSON):
    """The deck's repeating motif: a small hexagon, echoing the H3 grid."""
    return shape(slide, MSO_SHAPE.HEXAGON, x, y, size * 1.16, size,
                 fill=color, line=None)


def pointer(slide, x, y, w, text, size=12.0, color=CRIMSON, h=0.30,
            mark=True):
    """A required template pointer, rendered verbatim as a section heading."""
    if mark:
        hexmark(slide, x, y + 0.055, 0.125, color)
        x, w = x + 0.20, w - 0.20
    tb = txbox(slide, x, y, w, h)
    p = para(tb.text_frame, first=True, line=1.02)
    run(p, text, size, bold=True, color=color)
    return tb


def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w is not None:
        kw["width"] = E(w)
    if h is not None:
        kw["height"] = E(h)
    return slide.shapes.add_picture(path, E(x), E(y), **kw)


def delete_slide(prs, index):
    lst = prs.slides._sldIdLst
    sid = list(lst)[index]
    prs.part.drop_rel(sid.get(qn("r:id")))
    lst.remove(sid)


def retext(shape, runs):
    """Replace a shape's text, keeping its first run's formatting as the base."""
    tf = shape.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = runs
    return p0.runs[0]
