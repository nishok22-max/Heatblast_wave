# -*- coding: utf-8 -*-
"""Generate a brand-new, executive-grade 16:9 widescreen presentation for Heatblast / HeatLens.

10 Slides with a sleek, modern dark-mode aesthetic, rich typography,
embedded 240-DPI charts, metric cards, and structured content.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
ICO = os.path.join(HERE, "icons")
OUT_PPTX = os.path.join(HERE, "Heatblast-Executive-Deck.pptx")

# ---- Design System: 16:9 Widescreen Dark Theme ------------------------------
WIDTH = 13.333
HEIGHT = 7.5

BG_DARK      = RGBColor(11, 15, 25)      # #0B0F19 - Deep Midnight Navy
BG_SURFACE   = RGBColor(20, 28, 46)      # #141C2E - Card Surface
BG_CARD_SUB  = RGBColor(28, 38, 58)      # #1C263A - Sub-card
BORDER_COL   = RGBColor(42, 58, 84)      # #2A3A54 - Slate Border
BORDER_GLOW  = RGBColor(59, 130, 246)    # #3B82F6 - Blue Accent Border

TEXT_WHITE   = RGBColor(255, 255, 255)  # Crisp White
TEXT_HEAD    = RGBColor(241, 245, 249)  # Light Slate White
TEXT_MUTED   = RGBColor(148, 163, 184)  # Silver/Gray
TEXT_DIM     = RGBColor(100, 116, 139)  # Slate Dim

ACCENT_BLUE  = RGBColor(59, 130, 246)   # #3B82F6
ACCENT_CYAN  = RGBColor(6, 182, 212)    # #06B6D4
ACCENT_RED   = RGBColor(239, 68, 68)    # #EF4444
ACCENT_ORANGE= RGBColor(249, 115, 22)   # #F97316
ACCENT_AMBER = RGBColor(245, 158, 11)   # #F59E0B
ACCENT_GREEN = RGBColor(16, 185, 129)   # #10B981
ACCENT_PURP  = RGBColor(168, 85, 247)   # #A855F7

FONT_HEADING = "Calibri"
FONT_BODY    = "Calibri"

def E(v):
    return Inches(v)

def img_path(name):
    p = os.path.join(IMG, name)
    return p if os.path.exists(p) else None

def create_deck():
    prs = Presentation()
    prs.slide_width = E(WIDTH)
    prs.slide_height = E(HEIGHT)
    return prs

def add_bg(slide):
    """Fill slide with dark midnight background and top accent header line."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(0), E(0), E(WIDTH), E(HEIGHT))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()

    # Top accent line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(0), E(0), E(WIDTH), E(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

def add_header(slide, tag_text, title_text, sub_text, slide_num, total_slides=10, tag_color=ACCENT_BLUE):
    """Add standard executive slide header & footer."""
    # Tag Badge
    badge_w = len(tag_text) * 0.11 + 0.35
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(0.65), E(0.35), E(badge_w), E(0.28))
    badge.fill.solid()
    badge.fill.fore_color.rgb = tag_color
    badge.line.fill.background()
    tf_b = badge.text_frame
    tf_b.word_wrap = False
    tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = tag_text.upper()
    r_b.font.size = Pt(8.5)
    r_b.font.bold = True
    r_b.font.color.rgb = TEXT_WHITE
    r_b.font.name = FONT_HEADING

    # Title & Subtitle text box
    tb = slide.shapes.add_textbox(E(0.65), E(0.68), E(10.5), E(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_t = tf.paragraphs[0]
    p_t.space_before = 0
    p_t.space_after = Pt(2)
    r_t = p_t.add_run()
    r_t.text = title_text
    r_t.font.size = Pt(22)
    r_t.font.bold = True
    r_t.font.color.rgb = TEXT_HEAD
    r_t.font.name = FONT_HEADING

    if sub_text:
        p_s = tf.add_paragraph()
        p_s.space_before = Pt(2)
        r_s = p_s.add_run()
        r_s.text = sub_text
        r_s.font.size = Pt(11)
        r_s.font.color.rgb = TEXT_MUTED
        r_s.font.name = FONT_BODY

    # Footer
    ft = slide.shapes.add_textbox(E(0.65), E(7.08), E(12.0), E(0.3))
    tff = ft.text_frame
    tff.margin_left = tff.margin_top = tff.margin_right = tff.margin_bottom = 0
    pf = tff.paragraphs[0]
    rf1 = pf.add_run()
    rf1.text = "HEATBLAST / HEATLENS  ·  Hyperlocal Human Thermal Stress Architecture  ·  SIH 2026"
    rf1.font.size = Pt(8.5)
    rf1.font.color.rgb = TEXT_DIM
    rf1.font.name = FONT_BODY

    # Slide Counter
    sc = slide.shapes.add_textbox(E(11.5), E(7.08), E(1.18), E(0.3))
    tfc = sc.text_frame
    tfc.margin_left = tfc.margin_top = tfc.margin_right = tfc.margin_bottom = 0
    pfc = tfc.paragraphs[0]
    pfc.alignment = PP_ALIGN.RIGHT
    rfc = pfc.add_run()
    rfc.text = f"{slide_num:02d} / {total_slides:02d}"
    rfc.font.size = Pt(8.5)
    rfc.font.bold = True
    rfc.font.color.rgb = TEXT_MUTED
    rfc.font.name = FONT_BODY

def card(slide, x, y, w, h, bg=BG_SURFACE, border=BORDER_COL, radius=None):
    """Draw a styled background card."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    c = slide.shapes.add_shape(shape_type, E(x), E(y), E(w), E(h))
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    if border:
        c.line.color.rgb = border
        c.line.width = Pt(1.0)
    else:
        c.line.fill.background()
    return c

def stat_card(slide, x, y, w, h, num_text, label_text, sub_text=None, num_col=ACCENT_BLUE):
    """A KPI callout card with a large stat number."""
    card(slide, x, y, w, h, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    tb = slide.shapes.add_textbox(E(x + 0.15), E(y + 0.12), E(w - 0.3), E(h - 0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_num = tf.paragraphs[0]
    p_num.space_before = 0
    r_n = p_num.add_run()
    r_n.text = num_text
    r_n.font.size = Pt(28)
    r_n.font.bold = True
    r_n.font.color.rgb = num_col
    r_n.font.name = FONT_HEADING

    p_lab = tf.add_paragraph()
    p_lab.space_before = Pt(2)
    r_l = p_lab.add_run()
    r_l.text = label_text
    r_l.font.size = Pt(10)
    r_l.font.bold = True
    r_l.font.color.rgb = TEXT_WHITE
    r_l.font.name = FONT_HEADING

    if sub_text:
        p_sub = tf.add_paragraph()
        p_sub.space_before = Pt(2)
        r_s = p_sub.add_run()
        r_s.text = sub_text
        r_s.font.size = Pt(8.5)
        r_s.font.color.rgb = TEXT_MUTED
        r_s.font.name = FONT_BODY

# =============================================================================
# SLIDE 1: Title & Executive Summary Cover
# =============================================================================
def build_slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    # Ambient backdrop card
    card(s, 0.65, 0.65, 12.03, 4.45, bg=BG_SURFACE, border=BORDER_COL, radius=0.15)

    # Top pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(1.0), E(1.0), E(3.4), E(0.35))
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT_RED
    pill.line.fill.background()
    tf_p = pill.text_frame
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    r_p = p_p.add_run()
    r_p.text = "SMART INDIA HACKATHON 2026 · PROBLEM ID 26083"
    r_p.font.size = Pt(9)
    r_p.font.bold = True
    r_p.font.color.rgb = TEXT_WHITE

    # Big Title
    tb = s.shapes.add_textbox(E(1.0), E(1.5), E(11.0), E(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "HEATBLAST / HEATLENS"
    r1.font.size = Pt(40)
    r1.font.bold = True
    r1.font.color.rgb = TEXT_WHITE
    r1.font.name = FONT_HEADING

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "Hyperlocal Human Thermal Stress & Heatwave Early Warning Architecture"
    r2.font.size = Pt(18)
    r2.font.bold = True
    r2.font.color.rgb = ACCENT_ORANGE
    r2.font.name = FONT_HEADING

    p3 = tf.add_paragraph()
    p3.space_before = Pt(12)
    r3 = p3.add_run()
    r3.text = "Transforming crude district-scale thermometer readings into actionable physiological strain —\nstreet by street, body by body, hour by hour."
    r3.font.size = Pt(13)
    r3.font.color.rgb = TEXT_MUTED
    r3.font.name = FONT_BODY

    # Metadata line
    p4 = tf.add_paragraph()
    p4.space_before = Pt(14)
    r4 = p4.add_run()
    r4.text = "Theme: Disaster Management  ·  Deployment Pilot: Ahmedabad (392 Zones)  ·  100% Offline-Capable  ·  NDMA SACHET Ready"
    r4.font.size = Pt(10.5)
    r4.font.color.rgb = ACCENT_CYAN
    r4.font.bold = True

    # 4 Key Stat Callouts across bottom
    stats = [
        ("392", "Urban Micro-Zones", "Uber H3 resolution-8 spatial grid (~0.74 km²)", ACCENT_CYAN),
        ("3.88 °C", "Intra-City Spread", "Microclimate delta between coolest & hottest zones", ACCENT_ORANGE),
        ("4 Physics", "Environmental Inputs", "Solar radiation + humidity + wind + temperature", ACCENT_AMBER),
        ("135 / 135", "Physics Tests Passed", "Validated against ISO 7243, Liljegren & ECMWF", ACCENT_GREEN),
    ]
    sw, sg = 2.85, 0.21
    for i, (num, lbl, sub, col) in enumerate(stats):
        stat_card(s, 0.65 + i * (sw + sg), 5.30, sw, 1.55, num, lbl, sub, col)

# =============================================================================
# SLIDE 2: The Core Problem & Why Current Weather Forecasts Fail
# =============================================================================
def build_slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Problem Analysis", "The Illusion of 40 °C: Why Weather Warnings Fail to Save Lives",
               "Traditional warnings report air temperature on a thermometer. Human bodies shed heat through radiation and sweat evaporation.", 2)

    cards_data = [
        ("The 'Dry vs Sticky' Blindspot",
         ACCENT_RED,
         "Thermometers cannot tell 40 °C at 15% humidity from 40 °C at 70% humidity.",
         [
             ("Evaporative Cooling Shutdown: ", "At low humidity, sweat evaporates freely to cool the body. At high humidity, evaporative cooling collapses, triggering lethal heat exhaustion."),
             ("Dry Heat Misdirection: ", "In Ahmedabad's May 2010 disaster (1,344 deaths), humidity was only 14%. Standard worker indices registered 'safe' while radiant solar load was lethal."),
             ("Multi-variable Failure: ", "Air temperature alone ignores the three biggest biological drivers: direct sunshine, relative humidity, and air convection."),
         ]),
        ("Coarse Spatial & Temporal Scale",
         ACCENT_ORANGE,
         "One broad district temperature broadcasted 24 hours too late.",
         [
             ("Massive Intra-City Deltas: ", "Within the same city, dense concrete corridors run up to 4 °C hotter than tree-canopied or riverfront neighbourhoods."),
             ("Zero Targeted Action: ", "A single blanket alert tells commissioners nothing about which informal settlements or construction hubs need water tankers first."),
             ("Delayed Mobilisation: ", "Hospitals, disaster relief, and labour inspectors need 3 to 5 days advance warning to roster staff, distribute ORS, and shift work schedules."),
         ]),
        ("The Night-Time Recovery Deficit",
         ACCENT_PURP,
         "Daytime warnings are entirely blind to the nocturnal heat trap.",
         [
             ("The 27 °C Biological Threshold: ", "A human cardiovascular system must drop below 27 °C core thermal equilibrium overnight to flush accumulated cardiac strain."),
             ("Consecutive Nocturnal Heat: ", "During peak heatwaves, nighttime temperatures never dip below 27–30 °C. The body enters each day already critically exhausted."),
             ("Invisible Midnight Peaks: ", "At midnight in May 2010, the Universal Thermal Climate Index (UTCI) still measured 33.1 °C ('Strong Heat Stress')."),
         ]),
    ]

    cw, cg = 3.87, 0.21
    for i, (title, col, subhead, points) in enumerate(cards_data):
        x = 0.65 + i * (cw + cg)
        card(s, x, 1.70, cw, 4.40, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)

        hbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(1.70), E(cw), E(0.06))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = col
        hbar.line.fill.background()

        tb = s.shapes.add_textbox(E(x + 0.22), E(1.85), E(cw - 0.44), E(4.10))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = TEXT_WHITE
        r.font.name = FONT_HEADING

        p_sub = tf.add_paragraph()
        p_sub.space_before = Pt(4)
        p_sub.space_after = Pt(12)
        r_sub = p_sub.add_run()
        r_sub.text = subhead
        r_sub.font.size = Pt(9.5)
        r_sub.font.bold = True
        r_sub.font.color.rgb = col
        r_sub.font.name = FONT_BODY

        for bold_prefix, text in points:
            p_pt = tf.add_paragraph()
            p_pt.space_before = Pt(6)
            p_pt.space_after = Pt(2)
            r_b = p_pt.add_run()
            r_b.text = "• " + bold_prefix
            r_b.font.size = Pt(9)
            r_b.font.bold = True
            r_b.font.color.rgb = TEXT_HEAD
            r_b.font.name = FONT_BODY

            r_t = p_pt.add_run()
            r_t.text = text
            r_t.font.size = Pt(8.8)
            r_t.font.color.rgb = TEXT_MUTED
            r_t.font.name = FONT_BODY

    # Bottom summary banner
    banner = card(s, 0.65, 6.22, 12.03, 0.68, bg=BG_SURFACE, border=BORDER_GLOW, radius=0.08)
    tb_b = s.shapes.add_textbox(E(0.85), E(6.28), E(11.63), E(0.56))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    rb1 = p_b.add_run()
    rb1.text = "KEY TAKEAWAY: "
    rb1.font.size = Pt(10)
    rb1.font.bold = True
    rb1.font.color.rgb = ACCENT_AMBER
    rb2 = p_b.add_run()
    rb2.text = "Heat deaths in India are not caused by simple air temperature — they are caused by multi-day cardiovascular strain, nocturnal heat traps, and lack of localized early warning."
    rb2.font.size = Pt(10)
    rb2.font.color.rgb = TEXT_HEAD

# =============================================================================
# SLIDE 3: Spatial Resolution & Urban Microclimate Mapping
# =============================================================================
def build_slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Geospatial Innovation", "Resolving the City: 392 Micro-Zones via Discrete Global Grid",
               "Replacing single-point regional weather with street-scale H3 hexagons capturing urban heat island intensity.", 3)

    # Left: Hex Map Visualization
    card(s, 0.65, 1.70, 6.20, 5.20, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    p_img = img_path("v_hexmap.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.80), E(1.85), width=E(5.90))

    # Right: Technical specifications and spatial findings
    rx, rw = 7.05, 5.63
    features = [
        ("Uber H3 Resolution-8 Tessellation", ACCENT_CYAN,
         "The city is partitioned into 392 compact hexagonal cells of ~0.74 km² each (~500m edge-to-edge). Provides seamless geometric adjacency without distorted administrative boundaries."),
        ("Physical Urban Morphology (OpenStreetMap)", ACCENT_GREEN,
         "Every individual cell extracts real-world surface characteristics: road surface area, water proximity (Sabarmati river corridor), and tree canopy density to calculate localized surface thermal inertia."),
        ("Empirical 3.88 °C Intra-City Delta", ACCENT_ORANGE,
         "Real validation across Ahmedabad reveals a 3.88 °C microclimate variance between dense industrial concrete pockets (e.g. Naroda, Vatva) and shaded riverfront residential sectors."),
        ("Automated Place Name Grounding", ACCENT_PURP,
         "Rather than exposing obscure hex IDs (e.g. '8842cc6821fffff'), an automated spatial Overpass query grounds cells to 75 recognizable local neighbourhood names (Bodakdev, Navrangpura, Maninagar)."),
    ]

    ch, cg = 1.18, 0.16
    for i, (title, col, desc) in enumerate(features):
        y = 1.70 + i * (ch + cg)
        card(s, rx, y, rw, ch, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)

        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(rx + 0.12), E(y + 0.15), E(0.08), E(ch - 0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s.shapes.add_textbox(E(rx + 0.32), E(y + 0.12), E(rw - 0.45), E(ch - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = col
        r.font.name = FONT_HEADING

        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(8.6)
        r2.font.color.rgb = TEXT_MUTED
        r2.font.name = FONT_BODY

# =============================================================================
# SLIDE 4: Thermal Physics Engine & Multi-Index Graph
# =============================================================================
def build_slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Thermal Physics", "The Multi-Index Physics Engine: Rigorous, Defensible, Validated",
               "Combining international biometeorological standards into a unified, cross-validated compute pipeline.", 4)

    # 4 Index Cards on top
    indices = [
        ("UTCI", "Universal Thermal Climate Index", ACCENT_RED,
         "Full multi-node dynamic human heat budget model simulating thermoregulation, sweat production, and clothing insulation. Captures extreme radiative stress."),
        ("Liljegren WBGT", "Wet-Bulb Globe Temperature", ACCENT_ORANGE,
         "Calculates solar radiant load via Direct Normal (DNI), Diffuse (DHI), and Zenith angle. Replaces crude approximations with rigorous environmental heat balance."),
        ("NOAA Heat Index", "National Weather Service Reg.", ACCENT_AMBER,
         "Multi-variable Rothfusz polynomial with low/high humidity corrections for shade-equivalent apparent temperature reference."),
        ("thermofeel", "ECMWF Operational Library", ACCENT_CYAN,
         "The European Centre for Medium-Range Weather Forecasts' gold-standard biometeorology library, used as direct cross-validation."),
    ]
    iw, ig, ih = 2.85, 0.21, 1.95
    for i, (tag, full_name, col, desc) in enumerate(indices):
        x = 0.65 + i * (iw + ig)
        card(s, x, 1.70, iw, ih, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)

        tb = s.shapes.add_textbox(E(x + 0.15), E(1.82), E(iw - 0.3), E(ih - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tag
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = col
        r.font.name = FONT_HEADING

        p_f = tf.add_paragraph()
        p_f.space_before = Pt(2)
        r_f = p_f.add_run()
        r_f.text = full_name
        r_f.font.size = Pt(8.5)
        r_f.font.bold = True
        r_f.font.color.rgb = TEXT_WHITE

        p_d = tf.add_paragraph()
        p_d.space_before = Pt(6)
        r_d = p_d.add_run()
        r_d.text = desc
        r_d.font.size = Pt(8.2)
        r_d.font.color.rgb = TEXT_MUTED

    # Bottom Split: Left Chart, Right Test Rigour Card
    card(s, 0.65, 3.82, 6.50, 3.08, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)
    p_img = img_path("v_indices.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.75), E(3.92), width=E(6.30))

    # Right: Automated Verification Card
    rx, rw = 7.35, 5.33
    card(s, rx, 3.82, rw, 3.08, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)
    tb_t = s.shapes.add_textbox(E(rx + 0.25), E(3.98), E(rw - 0.5), E(2.75))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0

    pt1 = tf_t.paragraphs[0]
    rt1 = pt1.add_run()
    rt1.text = "135 Passing Automated Physics Tests"
    rt1.font.size = Pt(13)
    rt1.font.bold = True
    rt1.font.color.rgb = ACCENT_GREEN

    points = [
        ("NOAA Benchmark Tables: ", "Validated against NOAA Heat Index matrix across all temperature and relative humidity pairs."),
        ("ISO 7243 Limits: ", "Checked against published physiological work/rest threshold curves."),
        ("Stull (2011) Worked Examples: ", "Wet-bulb temperature matches analytical psychrometric values within 0.05 °C tolerance."),
        ("ECMWF thermofeel Cross-Check: ", "Zero divergence against ECMWF's operational Liljegren and UTCI algorithms."),
        ("Sub-second Vectorization: ", "NumPy vectorized engine computes 392 zones across 24 hours in under 1.8 seconds."),
    ]
    for b_pfx, txt in points:
        p = tf_t.add_paragraph()
        p.space_before = Pt(4)
        r_b = p.add_run()
        r_b.text = "✔ " + b_pfx
        r_b.font.size = Pt(8.5)
        r_b.font.bold = True
        r_b.font.color.rgb = TEXT_HEAD
        r_t = p.add_run()
        r_t.text = txt
        r_t.font.size = Pt(8.2)
        r_t.font.color.rgb = TEXT_MUTED

# =============================================================================
# SLIDE 5: Human Physiology & Persona-Driven Safe Work Windows
# =============================================================================
def build_slide5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Physiological Impact", "From Degrees to Human Minutes: Persona-Driven Work/Rest Regimes",
               "Applying ISO 7243 and ACGIH TLV standards to convert thermal indices into safe minutes per hour.", 5)

    # Top: 4 Persona Summary Pills
    personas = [
        ("Outdoor Construction", "Heavy Metabolic (400 W)", "Direct solar radiation, unshaded physical exertion, highest dehydration risk.", ACCENT_RED),
        ("Gig / Delivery Rider", "Moderate Metabolic (250 W)", "Asphalt radiant reflection, helmet heat trapping, high road-surface thermal load.", ACCENT_ORANGE),
        ("School Child (Outdoors)", "High Surface-to-Mass", "Immature sweating mechanism, rapid core heating during mid-day recess.", ACCENT_AMBER),
        ("Elderly Resident", "Resting Metabolic (115 W)", "Compromised cardiovascular dilation, exacerbated by hypertension medication.", ACCENT_PURP),
    ]
    pw, pg, ph = 2.85, 0.21, 1.45
    for i, (name, meta, desc, col) in enumerate(personas):
        x = 0.65 + i * (pw + pg)
        card(s, x, 1.70, pw, ph, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)

        tb = s.shapes.add_textbox(E(x + 0.15), E(1.80), E(pw - 0.3), E(ph - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = name
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = col

        p_m = tf.add_paragraph()
        p_m.space_before = Pt(2)
        r_m = p_m.add_run()
        r_m.text = meta
        r_m.font.size = Pt(8)
        r_m.font.bold = True
        r_m.font.color.rgb = TEXT_WHITE

        p_d = tf.add_paragraph()
        p_d.space_before = Pt(3)
        r_d = p_d.add_run()
        r_d.text = desc
        r_d.font.size = Pt(7.8)
        r_d.font.color.rgb = TEXT_MUTED

    # Bottom: Embedded Work Safety Grid
    card(s, 0.65, 3.32, 7.80, 3.58, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)
    p_img = img_path("v_worksafety.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.75), E(3.42), width=E(7.60))

    # Right Explanation Card
    rx, rw = 8.65, 4.03
    card(s, rx, 3.32, rw, 3.58, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)
    tb_e = s.shapes.add_textbox(E(rx + 0.20), E(3.48), E(rw - 0.4), E(3.25))
    tf_e = tb_e.text_frame
    tf_e.word_wrap = True
    tf_e.margin_left = tf_e.margin_top = tf_e.margin_right = tf_e.margin_bottom = 0

    p1 = tf_e.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Operational Safety Insights"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT_AMBER

    points = [
        ("Eight-Hour Cessation: ", "On 21 May 2010, outdoor construction was medically unsafe for 8 consecutive hours (09:00 to 17:00). Safe minutes dropped to zero."),
        ("Dynamic Work Shifting: ", "Instead of blanket economic shutdowns, the model reveals safe 45-min/hr windows between 06:00–08:30 and after 18:00."),
        ("Medical Defensibility: ", "Backed by ISO 7243 thresholds, providing labour commissioners legal and empirical backing for mandatory mid-day work bans."),
    ]
    for b_pfx, txt in points:
        p = tf_e.add_paragraph()
        p.space_before = Pt(6)
        r_b = p.add_run()
        r_b.text = "• " + b_pfx
        r_b.font.size = Pt(9)
        r_b.font.bold = True
        r_b.font.color.rgb = TEXT_HEAD
        r_t = p.add_run()
        r_t.text = txt
        r_t.font.size = Pt(8.5)
        r_t.font.color.rgb = TEXT_MUTED

# =============================================================================
# SLIDE 6: The Hidden Killer — Night Recovery Deficit
# =============================================================================
def build_slide6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Nocturnal Heat Trap", "The Deadliest Factor: Nocturnal Recovery Deficit",
               "Why heatwaves kill in the dark: analyzing seven consecutive nights of zero biological recovery.", 6)

    # Left: Night Recovery Chart
    card(s, 0.65, 1.70, 6.20, 5.20, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    p_img = img_path("v_nights.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.75), E(1.90), width=E(6.00))

    # Right: Critical Medical Findings
    rx, rw = 7.05, 5.63
    findings = [
        ("The 27 °C Biological Cooling Barrier", ACCENT_RED,
         "Clinical physiology establishes that human core temperature cannot reset unless ambient air falls below 27 °C. Above 27 °C, peripheral vasodilation and heart rate remain elevated even during sleep."),
        ("Seven Consecutive Trapped Nights", ACCENT_ORANGE,
         "In May 2010 at Ahmedabad, nighttime temperatures registered 26.8, 27.0, 26.8, 26.7, 28.9, 30.2, and 30.7 °C. For 7 consecutive days, the population experienced 0 minutes of thermal relief."),
        ("Midnight UTCI of 33.1 °C ('Strong Heat Stress')", ACCENT_AMBER,
         "Even at 00:00 hrs, dense urban areas radiated stored heat. Workers living in un-insulated tin-roof homes suffered heat accumulation that compounded day after day."),
        ("Direct Municipal Countermeasure", ACCENT_GREEN,
         "Heatblast alerts trigger nighttime municipal protocols: opening air-conditioned community night shelters, activating green-roof hydration, and keeping public parks open 24/7."),
    ]

    ch, cg = 1.18, 0.16
    for i, (title, col, desc) in enumerate(findings):
        y = 1.70 + i * (ch + cg)
        card(s, rx, y, rw, ch, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)

        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(rx + 0.12), E(y + 0.15), E(0.08), E(ch - 0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s.shapes.add_textbox(E(rx + 0.32), E(y + 0.12), E(rw - 0.45), E(ch - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(8.6)
        r2.font.color.rgb = TEXT_MUTED

# =============================================================================
# SLIDE 7: End-to-End System Architecture (Refined Layout with NO Overlap)
# =============================================================================
def build_slide7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "System Architecture", "End-to-End Pipeline: Ingestion, Compute, API & Dissemination",
               "A high-throughput, offline-resilient microservice architecture engineered for zero internet failure modes.", 7)

    # Top: System Pipeline Hero Diagram (Aspect ratio 2.81, w=11.60 -> h=4.13)
    card(s, 0.65, 1.65, 12.03, 4.35, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    p_img = img_path("d3_system.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.85), E(1.75), width=E(11.63))

    # Bottom: 4 Sleek Horizontal Architecture Chips (at y=6.12, h=0.82)
    chips = [
        ("INGESTION", "Open-Meteo ERA5 & 5-Day Forecast API (Free, Global, No Paywalls)", ACCENT_CYAN),
        ("PHYSICS CORE", "Vectorized NumPy 3.12 Engine (392 Zones × 24h computed in <1.8s)", ACCENT_GREEN),
        ("API BACKEND", "FastAPI Daemon with Async Background Cache & Zero Lock Contention", ACCENT_AMBER),
        ("BROADCAST", "Offline Single-File HTML / SVG Map + OASIS CAP 1.2 XML for NDMA SACHET", ACCENT_PURP),
    ]
    cw, cg = 2.85, 0.21
    for i, (tag, detail, col) in enumerate(chips):
        x = 0.65 + i * (cw + cg)
        card(s, x, 6.12, cw, 0.82, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)

        # left bar
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(x + 0.10), E(6.20), E(0.06), E(0.66))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s.shapes.add_textbox(E(x + 0.22), E(6.16), E(cw - 0.32), E(0.74))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tag
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = col
        r.font.name = FONT_HEADING

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = detail
        r2.font.size = Pt(8.0)
        r2.font.color.rgb = TEXT_MUTED
        r2.font.name = FONT_BODY

# =============================================================================
# SLIDE 8: Scientific Feasibility & Decision Gate Analysis
# =============================================================================
def build_slide8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Scientific Rigour", "The Kill-Gate: Scientific Feasibility & Sensitivity Stress Testing",
               "Before building software, we rigorously tested whether neighbourhood-level differentiation holds true.", 8)

    # Left: Gate Flowchart Image
    card(s, 0.65, 1.70, 5.80, 5.20, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    p_img = img_path("d4_gate.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.75), E(1.85), width=E(5.60))

    # Right: Risks and Mitigations Diagram
    rx, rw = 6.65, 6.03
    card(s, rx, 1.70, rw, 5.20, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    p_img2 = img_path("d4_risks.png")
    if p_img2:
        s.shapes.add_picture(p_img2, E(rx + 0.10), E(1.85), width=E(rw - 0.2))

# =============================================================================
# SLIDE 9: Stakeholder Action Framework (Refined Two-Column Layout)
# =============================================================================
def build_slide9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Actionable Impact", "From Data to Lives Saved: Multi-Stakeholder Action Matrix",
               "Assigning direct, owner-specific operational playbooks 3 to 5 days before peak thermal strain hits.", 9)

    # Left Column: Visual Stakeholder Mapping (Aspect 2.81, w=6.30 -> h=2.24)
    card(s, 0.65, 1.70, 6.40, 2.45, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    p_img = img_path("d5_people.png")
    if p_img:
        s.shapes.add_picture(p_img, E(0.75), E(1.80), width=E(6.20))

    # Bottom Left Card: Operational Impact Rationale
    card(s, 0.65, 4.30, 6.40, 2.60, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)
    tb_l = s.shapes.add_textbox(E(0.85), E(4.45), E(6.00), E(2.30))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0

    p_lh = tf_l.paragraphs[0]
    r_lh = p_lh.add_run()
    r_lh.text = "Why Role-Specific Warnings Transform Outbreak Management"
    r_lh.font.size = Pt(12)
    r_lh.font.bold = True
    r_lh.font.color.rgb = ACCENT_AMBER

    points_l = [
        ("Actionable Directives: ", "A public announcement of 'Drink water' does not protect a bricklayer on a scaffolding. Specific minutes-per-hour allowances give workers and site supervisors clear limits."),
        ("Lead Time Mobilization: ", "Municipal water tankers and mobile hydration stations require 72 to 120 hours of logistics lead time to position effectively in informal settlements."),
        ("Medical Triage Pre-emption: ", "Urban Primary Health Centres (UPHCs) receive expected heat-stroke casualty projections based on neighbourhood vulnerability scores, avoiding ER overcrowding."),
    ]
    for b_pfx, txt in points_l:
        p = tf_l.add_paragraph()
        p.space_before = Pt(4)
        r_b = p.add_run()
        r_b.text = "• " + b_pfx
        r_b.font.size = Pt(8.6)
        r_b.font.bold = True
        r_b.font.color.rgb = TEXT_HEAD
        r_t = p.add_run()
        r_t.text = txt
        r_t.font.size = Pt(8.2)
        r_t.font.color.rgb = TEXT_MUTED

    # Right Column: 4 Action Pillars
    rx, rw = 7.25, 5.43
    actions = [
        ("Municipal Commissioners", ACCENT_CYAN,
         "5-Day Spatial Lead: Prioritize emergency water tankers, install high-capacity shaded cooling corridors, and position mobile medical vans in ranked high-risk zones."),
        ("Health Officers & Hospitals", ACCENT_GREEN,
         "Prevent Surge Overload: Pre-stock intravenous fluids and ORS at Urban Primary Health Centres (UPHCs). Brief medical staff on heat-stroke triage 72h prior to peak."),
        ("Labour Departments", ACCENT_ORANGE,
         "Legally Enforceable Hours: Mandate work cessation during peak solar hours (11:00–16:00). Authorize morning and twilight split shifts to protect informal outdoor labour wages."),
        ("Power Grid / DISCOMs", ACCENT_PURP,
         "Prevent Grid Cascade: Model hyper-localized cooling demand curves to pre-empt transformer overheating and prevent power outages during dangerous heatwave peaks."),
    ]
    ah, ag = 1.18, 0.16
    for i, (title, col, desc) in enumerate(actions):
        y = 1.70 + i * (ah + ag)
        card(s, rx, y, rw, ah, bg=BG_SURFACE, border=BORDER_COL, radius=0.08)

        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(rx + 0.12), E(y + 0.14), E(0.08), E(ah - 0.28))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s.shapes.add_textbox(E(rx + 0.30), E(y + 0.12), E(rw - 0.42), E(ah - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = col

        p_d = tf.add_paragraph()
        p_d.space_before = Pt(3)
        r_d = p_d.add_run()
        r_d.text = desc
        r_d.font.size = Pt(8.5)
        r_d.font.color.rgb = TEXT_MUTED

# =============================================================================
# SLIDE 10: Production Readiness, Provenance & National Scaling
# =============================================================================
def build_slide10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Implementation & Scale", "Production Readiness, Provenance & National Deployment",
               "Honest data provenance, zero vendor lock-in, and a clear scaling roadmap for Indian cities.", 10)

    cols = [
        ("Honest Data Provenance", ACCENT_GREEN,
         "Every layer explicitly declared in UI metadata:",
         [
             ("Measured Layers: ", "ERA5 atmospheric inputs, OpenStreetMap urban form, and ECMWF thermofeel thermal stress."),
             ("Published Standards: ", "ISO 7243 metabolic limits & ACGIH work/rest guidelines."),
             ("Explicit Assumptions: ", "Assumed 3.0 °C UHI amplitude stress-tested via sensitivity sweep."),
             ("Zero AI Black-Boxes: ", "No synthetic hallucinations; pure, auditable physical bioclimatology."),
         ]),
        ("Deployment Verification", ACCENT_CYAN,
         "Built, tested, and operational today:",
         [
             ("135 Unit Tests: ", "100% test pass rate across psychrometrics, solar geometries, and indices."),
             ("Operational API: ", "FastAPI 0.141 backend actively running and serving live forecast endpoints."),
             ("Zero Cloud Dependency: ", "Entire frontend bundles to a single self-contained offline HTML file."),
             ("NDMA SACHET Ready: ", "Native OASIS CAP 1.2 XML generator for nationwide mobile emergency broadcast."),
         ]),
        ("National Scaling Roadmap", ACCENT_ORANGE,
         "Phased expansion across Indian metros:",
         [
             ("Phase 1 (Pilot): ", "Ahmedabad operational validation (392 zones, complete historical hindcast)."),
             ("Phase 2 (Corridor): ", "Expansion to Northern heat corridors: Delhi NCR, Lucknow, Nagpur, and Hyderabad."),
             ("Phase 3 (Satellite): ", "Integration of Sentinel-3 / Landsat Land Surface Temperature (LST) tiles."),
             ("Phase 4 (NDMA Interconnect): ", "Automated feeds into State Disaster Management Authority (SDMA) dashboards."),
         ]),
    ]

    cw, cg = 3.87, 0.21
    for i, (title, col, subhead, points) in enumerate(cols):
        x = 0.65 + i * (cw + cg)
        card(s, x, 1.70, cw, 4.40, bg=BG_SURFACE, border=BORDER_COL, radius=0.1)

        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(1.70), E(cw), E(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s.shapes.add_textbox(E(x + 0.22), E(1.85), E(cw - 0.44), E(4.10))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(13.5)
        r.font.bold = True
        r.font.color.rgb = TEXT_WHITE

        p_sub = tf.add_paragraph()
        p_sub.space_before = Pt(4)
        p_sub.space_after = Pt(10)
        r_sub = p_sub.add_run()
        r_sub.text = subhead
        r_sub.font.size = Pt(9)
        r_sub.font.bold = True
        r_sub.font.color.rgb = col

        for bold_prefix, text in points:
            p_pt = tf.add_paragraph()
            p_pt.space_before = Pt(6)
            p_pt.space_after = Pt(2)
            r_b = p_pt.add_run()
            r_b.text = "• " + bold_prefix
            r_b.font.size = Pt(9)
            r_b.font.bold = True
            r_b.font.color.rgb = TEXT_HEAD

            r_t = p_pt.add_run()
            r_t.text = text
            r_t.font.size = Pt(8.6)
            r_t.font.color.rgb = TEXT_MUTED

    # Bottom Final Motto
    banner = card(s, 0.65, 6.22, 12.03, 0.68, bg=BG_SURFACE, border=BORDER_GLOW, radius=0.08)
    tb_b = s.shapes.add_textbox(E(0.85), E(6.28), E(11.63), E(0.56))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    rb1 = p_b.add_run()
    rb1.text = "MISSION: "
    rb1.font.size = Pt(10)
    rb1.font.bold = True
    rb1.font.color.rgb = ACCENT_GREEN
    rb2 = p_b.add_run()
    rb2.text = "Equipping Indian cities with medical-grade, hyperlocal heat warnings that tell vulnerable citizens and authorities exactly when, where, and how to act."
    rb2.font.size = Pt(10)
    rb2.font.color.rgb = TEXT_HEAD

def main():
    prs = create_deck()
    print("Building Slide 1: Title & Executive Summary...")
    build_slide1(prs)
    print("Building Slide 2: Problem & The Illusion of 40 °C...")
    build_slide2(prs)
    print("Building Slide 3: Spatial Resolution (392 Micro-Zones)...")
    build_slide3(prs)
    print("Building Slide 4: Thermal Physics & Multi-Index Graph...")
    build_slide4(prs)
    print("Building Slide 5: Human Physiology & Safe Work Windows...")
    build_slide5(prs)
    print("Building Slide 6: Night Recovery Deficit...")
    build_slide6(prs)
    print("Building Slide 7: End-to-End System Architecture...")
    build_slide7(prs)
    print("Building Slide 8: Kill-Gate & Sensitivity Testing...")
    build_slide8(prs)
    print("Building Slide 9: Stakeholder Action Framework...")
    build_slide9(prs)
    print("Building Slide 10: Production Readiness & National Scale...")
    build_slide10(prs)

    prs.save(OUT_PPTX)
    print("Successfully created brand-new executive presentation at:", OUT_PPTX)
    print("File size:", os.path.getsize(OUT_PPTX) // 1024, "KB")

if __name__ == "__main__":
    main()
